#!/usr/bin/env python3
"""
export_to_pptx.py
Converts a slide-deck-creator JSON spec to a .pptx file using python-pptx.

Usage:
    python export_to_pptx.py deck.json [output.pptx]

Requires:
    pip install python-pptx

This is the Python alternative to generate_pptx.js. Use this if you prefer
a Python toolchain or need tighter control over python-pptx internals.
"""

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("Missing dependency. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_textbox(slide, text, x, y, w, h, **kwargs):
    """Add a text box and return the text frame for further styling."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text or ""
    p.alignment = {
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
    }.get(kwargs.get("align", "left"), PP_ALIGN.LEFT)

    run = p.runs[0] if p.runs else p.add_run()
    run.font.size    = Pt(kwargs.get("font_size", 18))
    run.font.bold    = kwargs.get("bold", False)
    run.font.italic  = kwargs.get("italic", False)
    run.font.color.rgb = hex_to_rgb(kwargs.get("color", "#FFFFFF"))
    run.font.name    = kwargs.get("font", "Calibri")
    return txBox


def set_bg(slide, color_hex: str):
    from pptx.oxml.ns import qn
    from lxml import etree
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def add_bullets(slide, bullets, x, y, w, h, color, font):
    if not bullets:
        return
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(15)
        run.font.color.rgb = hex_to_rgb(color)
        run.font.name = font
        p.space_after = Pt(6)


# --- Layout renderers ---

def render_title_hero(slide, spec, meta):
    add_textbox(slide, spec.get("headline", ""), 1, 1.8, 11.3, 1.5,
                font_size=44, bold=True, align="center",
                color="#FFFFFF", font=meta["font_heading"])
    add_textbox(slide, spec.get("subheadline", ""), 1.5, 3.5, 10, 0.8,
                font_size=20, align="center",
                color=spec.get("accent_color", meta["color_accent"]),
                font=meta["font_body"])


def render_split_text(slide, spec, meta):
    add_textbox(slide, spec.get("headline", ""), 0.4, 0.5, 5.5, 1.3,
                font_size=30, bold=True, color="#FFFFFF", font=meta["font_heading"])
    add_textbox(slide, spec.get("subheadline", ""), 0.4, 1.9, 5.5, 0.7,
                font_size=16, color=spec.get("accent_color", meta["color_accent"]),
                font=meta["font_body"])
    add_bullets(slide, spec.get("bullets", []), 0.4, 2.75, 5.5, 3.5,
                color="#CBD5E1", font=meta["font_body"])
    # Right panel placeholder
    tf_box = slide.shapes.add_textbox(Inches(6.5), Inches(0.5), Inches(6.4), Inches(6.3))
    tf_box.text_frame.text = f"[Visual: {spec.get('visual_idea', '')}]"
    tf_box.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    tf_box.text_frame.paragraphs[0].runs[0].font.color.rgb = hex_to_rgb("#94A3B8")
    tf_box.text_frame.paragraphs[0].runs[0].font.italic = True


def render_three_column(slide, spec, meta):
    add_textbox(slide, spec.get("headline", ""), 0.5, 0.3, 12.3, 1.0,
                font_size=32, bold=True, align="center",
                color="#FFFFFF", font=meta["font_heading"])
    add_textbox(slide, spec.get("subheadline", ""), 0.5, 1.4, 12.3, 0.55,
                font_size=14, align="center", color="#CBD5E1", font=meta["font_body"])
    columns = spec.get("columns", [])
    col_w = 3.8
    starts = [0.4, 4.75, 9.1]
    for i, col in enumerate(columns[:3]):
        x = starts[i]
        add_textbox(slide, col.get("heading", ""), x, 2.4, col_w, 0.55,
                    font_size=16, bold=True, align="center",
                    color="#FFFFFF", font=meta["font_heading"])
        add_textbox(slide, col.get("body", ""), x, 3.05, col_w, 2.5,
                    font_size=13, align="center", color="#CBD5E1", font=meta["font_body"])


def render_stat_grid(slide, spec, meta):
    add_textbox(slide, spec.get("headline", ""), 0.5, 0.25, 12.3, 0.9,
                font_size=28, bold=True, align="center",
                color="#FFFFFF", font=meta["font_heading"])
    add_textbox(slide, spec.get("subheadline", ""), 0.5, 1.2, 12.3, 0.5,
                font_size=13, align="center", color="#CBD5E1", font=meta["font_body"])
    stats = spec.get("stats", [])
    positions = [(0.5, 2.0), (6.9, 2.0), (0.5, 4.2), (6.9, 4.2)]
    accent = spec.get("accent_color", meta["color_accent"])
    for i, stat in enumerate(stats[:4]):
        x, y = positions[i]
        add_textbox(slide, stat.get("value", ""), x, y, 5.9, 1.1,
                    font_size=64, bold=True, align="center",
                    color=accent, font=meta["font_heading"])
        add_textbox(slide, stat.get("label", ""), x, y + 1.1, 5.9, 0.4,
                    font_size=16, align="center", color="#FFFFFF", font=meta["font_body"])
        add_textbox(slide, stat.get("context", ""), x, y + 1.55, 5.9, 0.3,
                    font_size=11, align="center", color="#94A3B8",
                    font=meta["font_body"], italic=True)


def render_quote(slide, spec, meta):
    accent = spec.get("accent_color", meta["color_accent"])
    add_textbox(slide, "\u201C", 0.3, 0.1, 2.0, 1.5,
                font_size=100, bold=True, color=accent, font=meta["font_heading"])
    add_textbox(slide, spec.get("headline", ""), 0.9, 1.3, 11.5, 2.2,
                font_size=34, bold=True, align="center",
                color="#FFFFFF", font=meta["font_heading"])
    add_textbox(slide, spec.get("subheadline", ""), 1.5, 3.7, 10.3, 0.6,
                font_size=15, align="center", color=accent,
                font=meta["font_body"], italic=True)


def render_cta(slide, spec, meta):
    accent = spec.get("accent_color", meta["color_accent"])
    add_textbox(slide, spec.get("headline", ""), 0.5, 1.2, 12.3, 1.3,
                font_size=38, bold=True, align="center",
                color="#FFFFFF", font=meta["font_heading"])
    add_textbox(slide, spec.get("subheadline", ""), 1.0, 2.65, 11.3, 0.7,
                font_size=16, align="center", color="#CBD5E1", font=meta["font_body"])
    cta = spec.get("cta", {})
    if cta.get("label"):
        add_textbox(slide, cta["label"], 4.2, 3.8, 4.9, 0.65,
                    font_size=18, bold=True, align="center",
                    color=meta["color_primary"], font=meta["font_heading"])
    if cta.get("url"):
        add_textbox(slide, cta["url"], 3.2, 4.65, 6.9, 0.3,
                    font_size=11, align="center", color="#94A3B8",
                    font=meta["font_body"], italic=True)


def render_full_bleed(slide, spec, meta):
    add_textbox(slide, f"[Image: {spec.get('visual_idea','')}]",
                0, 0, 13.3, 7.5,
                font_size=12, align="center", color="#94A3B8",
                font=meta["font_body"], italic=True)
    add_textbox(slide, spec.get("headline", ""), 0.5, 4.5, 12.3, 1.3,
                font_size=38, bold=True, align="center",
                color="#FFFFFF", font=meta["font_heading"])


LAYOUT_MAP = {
    "title-hero":        render_title_hero,
    "split-text-visual": render_split_text,
    "three-column":      render_three_column,
    "stat-grid":         render_stat_grid,
    "quote-callout":     render_quote,
    "cta-close":         render_cta,
    "full-bleed-image":  render_full_bleed,
}


def build_meta_defaults(raw: dict) -> dict:
    return {
        "color_primary":   raw.get("color_primary",   "#0F172A"),
        "color_secondary": raw.get("color_secondary", "#1E40AF"),
        "color_accent":    raw.get("color_accent",    "#38BDF8"),
        "font_heading":    raw.get("font_heading",    "Georgia"),
        "font_body":       raw.get("font_body",       "Calibri"),
    }


def build_pptx(deck: dict, output_path: str):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    raw_meta = deck.get("deck_meta", {})
    meta = build_meta_defaults(raw_meta)

    blank_layout = prs.slide_layouts[6]  # blank layout

    for spec in deck.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        bg_color = meta["color_primary"]
        set_bg(slide, bg_color)

        layout = spec.get("layout", "split-text-visual")
        renderer = LAYOUT_MAP.get(layout, render_split_text)
        renderer(slide, spec, meta)

        note = spec.get("speaker_note", "")
        if note:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = note

    prs.save(output_path)


def main():
    if len(sys.argv) < 2:
        print("Usage: python export_to_pptx.py <deck.json> [output.pptx]", file=sys.stderr)
        sys.exit(1)

    input_path  = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else Path(input_path).stem + ".pptx"

    with open(input_path, "r", encoding="utf-8") as f:
        deck = json.load(f)

    build_pptx(deck, output_path)

    n = len(deck.get("slides", []))
    print(f"✓ Written: {output_path}")
    print(f"  Slides: {n}")
    print(f"  Next: run the pptx skill's QA workflow to inspect visually.")
    print(f"  Visual QA: python -m markitdown {output_path}")


if __name__ == "__main__":
    main()
